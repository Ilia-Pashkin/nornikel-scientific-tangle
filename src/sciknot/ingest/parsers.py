"""Парсеры документов. Каждый возвращает список (label, text) — страница/слайд/лист."""

import logging
from pathlib import Path

log = logging.getLogger(__name__)

Unit = tuple[str, str]  # (метка источника, текст)


def parse_pdf(path: Path) -> list[Unit]:
    import fitz  # pymupdf

    units: list[Unit] = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                units.append((f"стр. {i}", text))
    return units


def parse_docx(path: Path) -> list[Unit]:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()
    return [("документ", text)] if text else []


def parse_pptx(path: Path) -> list[Unit]:
    from pptx import Presentation

    prs = Presentation(str(path))
    units: list[Unit] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
        if parts:
            units.append((f"слайд {i}", "\n".join(parts)))
    return units


def parse_xlsx(path: Path) -> list[Unit]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    units: list[Unit] = []
    for ws in wb.worksheets:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
            if len(rows) >= 500:  # защита от гигантских листов
                break
        if rows:
            units.append((f"лист «{ws.title}»", "\n".join(rows)))
    wb.close()
    return units


def parse_xls(path: Path) -> list[Unit]:
    """Legacy-формат Excel 97-2003 (xlrd 2.x поддерживает только .xls)."""
    import xlrd

    wb = xlrd.open_workbook(str(path))
    units: list[Unit] = []
    for sh in wb.sheets():
        rows: list[str] = []
        for r in range(min(sh.nrows, 500)):
            cells = [str(sh.cell_value(r, c)).strip()
                     for c in range(sh.ncols) if str(sh.cell_value(r, c)).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            units.append((f"лист «{sh.name}»", "\n".join(rows)))
    return units


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".docm": parse_docx,
    ".pptx": parse_pptx,
    ".xlsx": parse_xlsx,
    ".xls": parse_xls,
}

SKIPPED_EXTENSIONS = {".doc", ".zip", ".rar", ".001", ".002", ".gif"}


def parse_file(path: Path) -> list[Unit]:
    """Диспетчер. Возвращает [] для неподдерживаемых/битых файлов (с логом)."""
    parser = PARSERS.get(path.suffix.lower())
    if parser is None:
        if path.suffix.lower() not in SKIPPED_EXTENSIONS:
            log.warning("Неизвестное расширение, пропуск: %s", path)
        return []
    try:
        return parser(path)
    except Exception as e:  # битые файлы не должны ронять пайплайн
        log.warning("Ошибка парсинга %s: %s", path, e)
        return []
