"""Генерация презентации docs/presentation.pptx. Цифры тянутся живьём из Neo4j."""

import sys
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from sciknot.config import ROOT
from sciknot.graph.loader import get_driver

ACCENT = RGBColor(0x00, 0x57, 0xA6)
DARK = RGBColor(0x20, 0x24, 0x2A)
GRAY = RGBColor(0x5A, 0x62, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def stats() -> dict:
    d = get_driver()
    try:
        with d.session() as s:
            out = {}
            for lbl in ("Document", "Chunk", "Material", "Process", "Equipment", "Experiment", "Facility", "Topic", "Expert", "Parameter"):
                out[lbl] = s.run(f"MATCH (n:{lbl}) RETURN count(n) AS c").single()["c"]
            out["Relations"] = s.run("MATCH ()-[r:RELATES]->() RETURN count(r) AS c").single()["c"]
            out["Contradicts"] = s.run(
                "MATCH ()-[r:RELATES {type:'contradicts'}]->() RETURN count(r) AS c").single()["c"]
            out["Affiliations"] = s.run(
                "MATCH ()-[r:AFFILIATED_WITH]->() RETURN count(r) AS c").single()["c"]
        return out
    except Exception:
        return {k: "—" for k in ("Document", "Chunk", "Material", "Process", "Equipment", "Experiment",
                                 "Facility", "Topic", "Expert", "Parameter", "Relations",
                                 "Contradicts", "Affiliations")}
    finally:
        d.close()


def add_slide(prs, title: str, bullets: list, accent_first=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), SLIDE_W - Inches(1.2), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), SLIDE_W - Inches(1.4), SLIDE_H - Inches(2.1))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        level = 0
        text = item
        if isinstance(item, tuple):
            level, text = item
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.level = level
        para.font.size = Pt(20 if level == 0 else 17)
        para.font.color.rgb = DARK if level == 0 else GRAY
        para.space_after = Pt(10)
        if accent_first and i == 0:
            para.font.bold = True
            para.font.color.rgb = ACCENT
    return slide


def main():
    st = stats()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Титул
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), SLIDE_W - Inches(1.6), Inches(2.6))
    tf = bg.text_frame
    p = tf.paragraphs[0]
    p.text = "Научный клубок"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = "Единая карта знаний R&D для горно-металлургической отрасли"
    p2.font.size = Pt(24)
    p2.font.color.rgb = DARK
    p3 = tf.add_paragraph()
    p3.text = "Граф знаний + семантический поиск (GraphRAG) · Хакатон Норникель · июль 2026"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRAY

    add_slide(prs, "Проблема: знания есть — структуры нет", [
        "Институциональная память рассеяна по отчётам, презентациям и личным архивам",
        "Команды дублируют литобзоры, потому что не видят уже выполненную работу",
        "Междисциплинарный поиск невозможен: данные разрознены по форматам и папкам",
        "Ответ на инженерный вопрос = дни ручного сбора из десятков источников",
        "Нет единой верифицированной базы → противоречивые выводы и споры интерпретаций",
    ])

    add_slide(prs, "Решение: GraphRAG-карта знаний", [
        "PDF / DOCX / PPTX / XLSX / XLS → парсинг и чанкинг (метаданные: категория, журнал, год, страница)",
        "Сканы и изображения → VLM (облако Gemini / локально Qwen+mmproj) → текст в общий индекс",
        "LLM-экстракция: материалы, процессы, оборудование, эксперименты, организации, параметры, эксперты, связи",
        "Neo4j: граф знаний + векторный индекс (bge-m3) + fulltext — в одной базе",
        "Запрос на естественном языке → разбор → векторный + fulltext поиск + обход графа",
        "Синтез ответа со ссылками на источники [N], аналитикой и оценкой достоверности",
    ], accent_first=False)

    add_slide(prs, "Проиндексированный корпус (цифры живые)", [
        f"Документов: {st['Document']}   ·   Чанков: {st['Chunk']}",
        f"Материалы и вещества: {st['Material']}   ·   Темы: {st['Topic']}",
        f"Процессы: {st['Process']}   ·   Оборудование: {st['Equipment']}   ·   Эксперименты: {st['Experiment']}",
        f"Эксперты (люди): {st['Expert']}   ·   Организации: {st['Facility']}   ·   связей «работает в»: {st['Affiliations']}",
        f"Параметры с числами и единицами: {st['Parameter']}",
        f"Связей в графе: {st['Relations']}, из них противоречий (contradicts): {st['Contradicts']}",
        (1, "Проиндексированы Обзоры, Статьи, Доклады + 8 выпусков журналов; остальные журналы и конференции — доиндексация одной командой"),
    ], accent_first=True)

    add_slide(prs, "Онтология и верификация знаний", [
        "Типы сущностей: Material · Process · Equipment · Experiment · Facility · Parameter · Expert · Topic · Document",
        "Связи: uses_material, produces_output, operates_at_condition, uses_equipment, applied_for, expert_in, contradicts, validated_by",
        "Верификация: confidence связи = число независимых источников; created_at / updated_at — дата актуализации факта",
        "Эксперты отделены от организаций: Expert — люди, Facility — заводы и лаборатории, связь AFFILIATED_WITH",
        "Каждый факт трассируется до документа, страницы и года",
        "Параметры хранят значение + единицу + оператор: «сульфаты ≤ 300 мг/л» — это данные, а не текст",
    ])

    add_slide(prs, "Пример: запрос из ТЗ (проверено по корпусу)", [
        "«Какие технические решения организации циркуляции католита при электроэкстракции никеля описаны в мировой практике, и какая скорость потока считается оптимальной?»",
        (1, "Ответ сгруппирован по технологиям, все цифры сверены с документами:"),
        (1, "— диафрагменные ячейки: циркуляция католита 20–30 л/ч [11]"),
        (1, "— анодные мешки: 1,5–5 м³/ч в пилотных испытаниях Outotec / Norilsk Nickel [12]"),
        (1, "— промышленные заводы Niihama и Nikkelverk: 0,035–0,07 м³/мин [14] + ячейки EMEW"),
        (1, "— честная оговорка: «оптимальность» в источниках не доказана — приведены рабочие диапазоны"),
        "Анти-галлюцинация: на вопрос о закачке шахтных вод (темы нет в корпусе) система отвечает «данных нет» и предлагает смежное — вместо выдумки",
    ], accent_first=True)

    add_slide(prs, "Многопараметрические запросы", [
        "Материал + процесс + условия + география + время — в одном вопросе",
        "Числовые диапазоны: параметры извлечены как число + единица + оператор",
        "География: автоматическая разметка RU / зарубежная практика / обе",
        "Временные рамки: «за последние 5 лет» → фильтр по году документа",
        "Мультиязычность: терминология RU/EN сопоставляется при разборе запроса и поиске",
    ])

    add_slide(prs, "Аналитика в каждом ответе — автоматически", [
        "Обзор источников: группировка по методу, году, географии, уровню детализации",
        "Консенсус и разногласия: что подтверждено несколькими источниками, где данные расходятся",
        "Достоверность: оценка уверенности + число подтверждающих источников",
        "Пробелы знаний: неизученные комбинации «материал–режим–условие», темы только RU / только зарубеж",
        "Рекомендации: похожие кейсы из смежных областей, эксперты и организации, смежные темы",
        "Сравнительные вопросы («А vs Б», «RU vs мир») → таблица сравнения; отчёт — в Markdown одной кнопкой",
        "Вкладки диалогов: несколько запросов параллельно, история переживает перезагрузку",
    ])

    add_slide(prs, "3D-граф знаний (стиль Obsidian)", [
        "Цепочки «материал → процесс → оборудование → результат» через типизированные связи; результат — двухцветная полусфера",
        f"Противоречия — красно-белые рёбра-«зебры» ({st['Contradicts']} в корпусе): конфликтующие данные видны сразу",
        "⚠-подсветка пробелов: сущности с покрытием ≤2 документов — кандидаты на НИР",
        "Эксперты и организации по теме запроса — оверлей вокруг любой сущности",
        "Hover по связи: тип, число подтверждающих источников, дата актуализации факта",
    ])

    add_slide(prs, "Работа с данными", [
        "Загрузка PDF/DOCX/PPTX/XLSX/XLS через UI с прогрессом всех стадий индексации",
        "Дедупликация по содержимому: тот же файл под другим именем не задублирует граф",
        "Сканы распознаются VLM автоматически; галочка «Индексировать изображения» — все графики и схемы в индекс",
        "К поисковому запросу можно приложить картинку: кнопка, drag-and-drop или Ctrl+V",
        "Блок «Пробелы в знаниях»: слабо освещённые процессы и геоперекосы тем по всему корпусу",
        "Полная очистка данных одной кнопкой (двойное подтверждение)",
    ])

    add_slide(prs, "Локальный контур и управление моделями", [
        "Полностью офлайн-режим: llama.cpp + GGUF — данные не покидают контур компании",
        "Переключение облако ↔ локальная модель в один клик; при уходе в облако локальная модель гасится автоматически (VRAM)",
        "Слоты LLM и Эмбеддинг: фиксированные порты, автозапуск сервера при первом запросе",
        "Vision: локально mmproj, в облаке Gemini — сканы при индексации, изображения-вложения к запросу",
        "Кнопка «Стоп»: честный обрыв генерации; журнал действий (аудит)",
        "Ключи API и данные — вне репозитория (.env, .gitignore)",
    ])

    add_slide(prs, "Масштабируемость", [
        "Инкрементальный пайплайн: чекпоинты по chunk_id — дозагрузка без пересборки",
        "Новая категория данных = одна команда: run_ingest.py --categories «Журналы»",
        "Новый домен = расширение онтологии в промпте экстракции (без изменения кода)",
        "Хранилище: Neo4j держит миллионы узлов; поиск — векторный индекс + fulltext",
        "LLM-слой заменяем: любой OpenAI-совместимый API (DeepSeek / GPT / Claude / локальная)",
    ])

    add_slide(prs, "Стек и результат", [
        "Python · uv · Neo4j 5 (граф + векторы + fulltext) · FastAPI + ванильный JS (тёмная тема, SSE-стриминг)",
        "LLM: DeepSeek V4 (flash — экстракция, pro — синтез) · vision Gemini 3.1 Flash Lite · эмбеддинги bge-m3 · локально llama.cpp",
        "Полный цикл: сырые документы → верифицированный граф → ответ с источниками",
        "Все 4 эталонных запроса из ТЗ проверены против корпуса: цифры ответов сходятся с документами дословно, на отсутствующую тему — честное «данных нет»",
    ], accent_first=False)

    out = ROOT / "docs" / "presentation.pptx"
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"OK -> {out}")


if __name__ == "__main__":
    main()
